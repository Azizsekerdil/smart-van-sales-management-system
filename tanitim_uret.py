# -*- coding: utf-8 -*-
"""
tanitim_uret.py — Van Sales tanıtım sunumunu tek kaynaktan üretir
═══════════════════════════════════════════════════════════════════════════════
Üretilen dosyalar (docs/sunum/ altına):

    Van_Sales_Tanitim.pptx        ← Türkçe, koyu tema (ekran/projeksiyon)
    Van_Sales_Tanitim.pdf
    Van_Sales_Tanitim.html        ← slayt PNG'leri gömülü mobil sunum
    Van_Sales_Tanitim_Baski.pptx  ← beyaz zemin, mono yazıcı için
    Van_Sales_Tanitim_Baski.pdf

    Van_Sales_Intro_EN.pptx / .pdf / .html
    Van_Sales_Intro_EN_Print.pptx / .pdf

Neden bu betik var
------------------
Beş dosyayı elle düzenlemek beş ayrı iş demek; üstelik PDF ve HTML PPTX'ten
türetildiği için elle bakım tutarsızlığı garanti eder. Betik zinciri kurar:
slaytları veriden üretir, PowerPoint'e PDF ve PNG ürettirir, HTML'i o PNG'lerle
yazar. Tek kaynak SLAYTLAR sözlüğüdür.

Tasarım kararı — neden şekilleri koddan kuruyoruz
------------------------------------------------
Finansal Analiz Pro sunumunda mevcut bir şablon slaydı derin kopyalanıyordu;
orada bu doğru karardı çünkü elle tasarlanmış bir sunum zaten vardı. Burada yok.
Şablon kopyalamak için önce elle bir şablon yapmak gerekirdi ve o şablon da bu
betiğin dışında, bakımsız bir dosya olarak kalırdı. Renk/ölçü sabitlerini koda
almak, tasarımı da sürüm kontrolüne sokar.

Yeniden çalıştırılabilirlik
---------------------------
Üretim her koşuda sıfırdan başlar (Presentation() boş sunum). Var olan çıktılar
üzerine yazılır. İdempotency ispatlanacak bir şey değil, tanım gereği sağlanır.

Kullanım:
    python tanitim_uret.py                 # hepsini üret (TR + EN)
    python tanitim_uret.py --dil tr        # yalnız Türkçe
    python tanitim_uret.py --sadece-pptx   # PowerPoint'e dokunma
    python tanitim_uret.py --kontrol       # üretmeden mevcut durumu raporla

NOT: PDF/PNG/HTML üretimi Microsoft PowerPoint gerektirir (COM otomasyonu).
     PowerPoint yoksa PPTX yine üretilir, diğerleri atlanır.

Çıkış kodu 0 = başarılı.
"""
import argparse
import base64
import glob
import os
import re
import shutil
import sys
import tempfile

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass


def _iz(isaret):
    """Konsol UTF-8 desteklemiyorsa ASCII yedeğe düş."""
    yedek = {"✓": "[OK]", "✗": "[!]", "→": "->", "·": "-", "…": "..."}
    try:
        isaret.encode(sys.stdout.encoding or "utf-8")
        return isaret
    except Exception:
        return yedek.get(isaret, "?")


KOK = os.path.dirname(os.path.abspath(__file__))
# Yayımlanan sürüm ayrı bir dizine ve ayrı bir ada üretilir. Özgün dosyanın
# üzerine yazmak, "hangi PDF denetimden geçti" sorusunu cevaplanamaz hâle
# getirir; iki dosya yan yana durur ve karşılaştırılabilir.
CIKTI = os.path.join(KOK, "docs", "presentation")
EKRANLAR = os.path.join(CIKTI, "ekranlar")
DILLER = ("tr", "en")


# ══════════════════════════════════════════════════════════════════════════
# SÜRÜM — koddan okunur, elle güncellenmez
# ══════════════════════════════════════════════════════════════════════════
def _surum():
    """(sürüm, kaynak) — CHANGELOG'un en üstteki sürüm başlığından."""
    yol = os.path.join(KOK, "CHANGELOG.md")
    try:
        with open(yol, encoding="utf-8") as f:
            for satir in f:
                m = re.match(r"^##\s*\[([0-9]+\.[0-9]+\.[0-9]+)\]", satir.strip())
                if m:
                    return "v" + m.group(1), "CHANGELOG.md"
    except Exception:
        pass
    return "v1.0.0", "varsayılan"


# ══════════════════════════════════════════════════════════════════════════
# ÖLÇÜM — slaytlardaki her sayı koddan ölçülür
# ══════════════════════════════════════════════════════════════════════════
# Sunumda "71 tablo, 259 uç nokta" gibi elle yazılmış rakamlar vardı. Bir
# sunum, kaynağıyla birlikte yaşlanmayan tek belgedir: kod büyür, slayt
# büyümez ve rakam sessizce yanlışa döner. Buradaki değerler her üretimde
# uygulamanın kendisinden okunur.
#
# Ölçüm başarısızsa slayt üretilir ama sayı yerine "—" konur. Doğrulanamayan
# bir rakamı basmaktansa boş bırakmak yeğdir.
def _olcum():
    """Uygulamayı içe aktarıp ölçülebilir her şeyi say."""
    import glob as _glob
    import json as _json
    import tempfile as _tempfile

    sonuc = {}
    onceki = dict(os.environ)
    try:
        gecici = _tempfile.mkdtemp(prefix="vs_olcum_")
        os.environ.update({
            "VS_ENV": "test",
            "VS_DATABASE_URL": "sqlite:///" + os.path.join(gecici, "olcum.db").replace("\\", "/"),
            "VS_SECRET_KEY": "o" * 70,
            "VS_LOG_DIR": os.path.join(gecici, "logs"),
            "VS_BACKUP_DIR": os.path.join(gecici, "backups"),
            "VS_LMSTUDIO_ENABLED": "false",
            "VS_NVIDIA_ENABLED": "false",
            "VS_CLAUDE_ENABLED": "false",
        })
        backend = os.path.join(KOK, "backend")
        if backend not in sys.path:
            sys.path.insert(0, backend)

        # ``import app.models`` rebinds the name ``app`` in this scope, so the
        # FastAPI instance is aliased on import; otherwise the route walk below
        # runs against the package and finds nothing.
        from fastapi.routing import APIRoute            # noqa: PLC0415
        from app.main import app as fastapi_app         # noqa: PLC0415
        from app.models.base import Base                # noqa: PLC0415
        import app.models                               # noqa: F401,PLC0415
        import app.compliance.models                    # noqa: F401,PLC0415
        from app.core.permissions import (              # noqa: PLC0415
            RESOURCES, ROLES, permission_code,
        )

        # ``app.routes`` düz bir liste DEĞİLDİR. FastAPI 0.141'den beri
        # ``include_router`` alt yönlendiriciyi olduğu gibi saklar; üst listede
        # yalnızca üç yol görünür. Basit bir ``isinstance(r, APIRoute)`` süzgeci
        # hata vermez — sessizce 3 döner ve slayta yanlış rakam yazar. Bu yüzden
        # ağaç geziliyor ve sonuç sıfırsa ölçüm başarısız sayılıyor.
        def _uc_noktalar(dugum, gorulen=None):
            gorulen = set() if gorulen is None else gorulen
            if dugum is None or id(dugum) in gorulen:
                return []
            gorulen.add(id(dugum))
            baglamlar = getattr(dugum, "effective_route_contexts", None)
            if callable(baglamlar):
                return [
                    (c.path, set(getattr(c, "methods", ()) or ()))
                    for c in baglamlar()
                ]
            bulunan = []
            for r in getattr(dugum, "routes", []) or []:
                if isinstance(r, APIRoute):
                    bulunan.append((r.path, set(r.methods or ())))
                bulunan += _uc_noktalar(r, gorulen)
                bulunan += _uc_noktalar(getattr(r, "app", None), gorulen)
            return bulunan

        uclar = 0
        for yol, yontemler in _uc_noktalar(fastapi_app):
            if yol.startswith("/api/v1"):
                uclar += len(yontemler - {"HEAD", "OPTIONS"})
        if uclar == 0:
            raise RuntimeError("uc nokta sayilamadi")
        sonuc["uc_nokta"] = uclar

        md = Base.metadata
        sonuc["tablo"] = len(md.tables)
        indeksler = set()
        for t in md.tables.values():
            for i in t.indexes:
                indeksler.add((t.name, i.name))
            for c in t.columns:
                if c.index:
                    indeksler.add((t.name, "ix_%s_%s" % (t.name, c.name)))
                if c.unique:
                    indeksler.add((t.name, "uq_%s_%s" % (t.name, c.name)))
        sonuc["indeks"] = len(indeksler)

        sonuc["izin"] = len({permission_code(r.key, a) for r in RESOURCES for a in r.actions})
        sonuc["rol"] = len(ROLES)

        sonuc["servis"] = len([
            y for y in _glob.glob(os.path.join(backend, "app", "services", "*.py"))
            if not y.endswith("__init__.py")
        ]) + len([
            y for y in _glob.glob(os.path.join(backend, "app", "compliance", "services", "*.py"))
            if not y.endswith("__init__.py")
        ])

        try:
            from app.reports.engine import REPORTS      # noqa: PLC0415
            sonuc["rapor"] = len(REPORTS)
        except Exception:
            pass

        app_tsx = os.path.join(KOK, "frontend", "src", "App.tsx")
        with open(app_tsx, encoding="utf-8") as f:
            metin = f.read()
        sonuc["ekran"] = len({
            y for y in re.findall(r"\{\s*path:\s*'([^']+)'", metin) if ":" not in y
        })

        with open(os.path.join(KOK, "frontend", "src", "locales", "tr.json"), encoding="utf-8") as f:
            _json.load(f)
        sonuc["dil"] = len(DILLER)
    except Exception as exc:                            # pragma: no cover
        print("   %s ölçüm yapılamadı: %s" % (_iz("✗"), str(exc)[:120]))
    finally:
        os.environ.clear()
        os.environ.update(onceki)
    return sonuc


OLCUM = _olcum()


def _s(anahtar):
    """Ölçülen değer, ölçülemediyse '—'."""
    deger = OLCUM.get(anahtar)
    return "—" if deger in (None, 0) else str(deger)


# ══════════════════════════════════════════════════════════════════════════
# PALET
# ══════════════════════════════════════════════════════════════════════════
# Uygulamanın kendi Tailwind paletiyle aynı: shell-900 zemin, brand-600 vurgu.
# Sunum ile ürün aynı görünsün diye değerler birebir kopyalandı.
Z = {
    "zemin":    "0F172A",   # shell-900
    "zemin2":   "111C33",   # kapak degradesi için biraz daha açık
    "kart":     "1E293B",   # shell-800
    "kart_cizgi": "334155",  # shell-700
    "marka":    "4F46E5",   # brand-600
    "marka_ac": "818CF8",   # brand-400
    "vurgu":    "2DD4BF",   # teal — sayılar ve rozetler
    "yazi":     "F1F5F9",
    "yazi_ac":  "FFFFFF",
    "sonuk":    "94A3B8",   # shell-400
    "sonuk2":   "64748B",   # shell-500
    "ok":       "10B981",
    "uyari":    "F59E0B",
    "hata":     "EF4444",
}

# Serbest lisanslı bir yazı tipi seçilir; sunum kamuya açık bir depoda
# yayımlanıyor ve tescilli bir yazı tipini gömülü olarak dağıtmak ayrı bir
# lisans sorusu doğurur. DejaVu Sans (Bitstream Vera türevi) Türkçe
# karakterlerin tamamını taşır.
#
# UYARI: PowerPoint yalnızca KURULU yazı tiplerini gömer. DejaVu Sans kurulu
# değilse sessizce başka bir yazı tipiyle değiştirir ve PDF onu gömer — yani
# bu satırı değiştirmek tek başına yetmez. Üretimden sonra
# ``THIRD_PARTY_NOTICES.md`` içindeki kayıt, PDF'te gerçekten gömülü olan
# yazı tipiyle doğrulanmalıdır.
YAZI_TIPI = "DejaVu Sans"


# ══════════════════════════════════════════════════════════════════════════
# BASKI PALETİ — mono lazer yazıcı için
# ══════════════════════════════════════════════════════════════════════════
# Koyu zemin A4'te sayfa başına ~%95 toner demek. Baskı sürümü zemini beyaza,
# yazıyı siyaha çevirir; vurgu renkleri koyu griye iner ki gri tonlamada
# birbirinden ayrılabilsinler.
BASKI_ZEMIN = {
    Z["zemin"]:      "FFFFFF",
    Z["zemin2"]:     "FFFFFF",
    Z["kart"]:       "F4F6F9",
    Z["kart_cizgi"]: "C7D0DA",
    Z["marka"]:      "2F3A8C",
    Z["marka_ac"]:   "5A64B8",
    Z["vurgu"]:      "0F766E",
    Z["ok"]:         "15803D",
    Z["uyari"]:      "92400E",
    Z["hata"]:       "B91C1C",
}
BASKI_YAZI = {
    Z["yazi"]:     "111827",
    Z["yazi_ac"]:  "000000",
    Z["sonuk"]:    "4B5563",
    Z["sonuk2"]:   "6B7280",
    Z["vurgu"]:    "0F766E",
    Z["marka_ac"]: "2F3A8C",
    Z["ok"]:       "15803D",
    Z["uyari"]:    "92400E",
    Z["hata"]:     "B91C1C",
}


# ══════════════════════════════════════════════════════════════════════════
# İÇERİK — tek kaynak
# ══════════════════════════════════════════════════════════════════════════
# Buradaki her sayı programdan ölçülmüştür (bkz. IMPLEMENTATION_PLAN.md).
# Ölçülmemiş bir sayı sunuma girmez.

SLAYTLAR_TR = [
    {"tip": "kapak",
     "baslik": "Akıllı Sıcak Satış\nYönetim Sistemi",
     "alt": "Yiyecek & içecek dağıtımı için uçtan uca van sales platformu",
     "etiketler": ["Sıcak Satış", "FEFO Stok", "Rota Optimizasyonu", "Yerel Yapay Zeka"],
     "dip": "Fabrikadan bakkal rafına — tek sistem"},

    {"tip": "vurgu",
     "ust": "SORUN",
     "baslik": "Sıcak satışta gün sonunda\nherkes farklı rakam söyler.",
     "alt": "Araçtaki stok, kesilen fatura, tahsil edilen nakit ve fiziki sayım "
            "ayrı ayrı tutulduğunda; akşam çıkan farkın kaynağı bulunamaz."},

    {"tip": "kartlar",
     "baslik": "Sahada kırılan dört nokta",
     "alt": "Van sales operasyonunda en sık maliyet üreten kalemler",
     "kartlar": [
         ("Araç stoğu tahmini", "Sabah ne yüklendiği not defterinde, gün içinde "
          "ne satıldığı fişte. İkisi akşam tutmayınca fark kimsenin üstüne yazılamaz."),
         ("SKT kaçağı", "Son kullanma tarihi yakın ürün rafta beklerken yeni parti "
          "satılır. Fire, ancak imha edilirken fark edilir."),
         ("Rota körlüğü", "Rota alışkanlıkla kurulur. Kat edilen kilometre ile "
          "yapılan ciro arasındaki ilişki ölçülmez."),
         ("Tahsilat riski", "Kredi limiti kâğıt üstündedir; plasiyer sahada limiti "
          "aşan satışı yaptıktan sonra öğrenilir."),
     ]},

    {"tip": "zincir",
     "baslik": "Sistemin kapsadığı zincir",
     "alt": "Her halka bir öncekinden veri alır — hiçbir yerde elle aktarım yok",
     "adimlar": ["Fabrika", "Merkez Depo", "Bölge Deposu", "Araç Yükleme",
                 "Satış Aracı", "Plasiyer", "Günlük Rota", "Müşteri",
                 "Sıcak Satış", "Fatura / İrsaliye", "Tahsilat", "İade",
                 "Araç Sayımı", "Gün Sonu Mutabakatı"]},

    {"tip": "sayilar",
     "baslik": "Sistemin ölçüsü",
     "alt": "Rakamlar çalışan sistemden ölçüldü, hedef değil",
     "sayilar": [(_s("tablo"), "veritabanı tablosu"), (_s("uc_nokta"), "API uç noktası"),
                 (_s("ekran"), "arayüz ekranı"), (_s("izin"), "yetki tanımı"),
                 (_s("servis"), "iş servisi"), (_s("rapor"), "hazır rapor"),
                 (_s("indeks"), "veritabanı indeksi"), (_s("dil"), "dil (TR / EN)")]},

    {"tip": "ekran_tam", "dosya": "02-panel",
     "baslik": "Kontrol paneli",
     "alt": "20 KPI ve 7 canlı grafik — panel ile raporlar aynı toplulaştırmadan beslenir"},

    {"tip": "bolum", "no": "01", "baslik": "Sıcak Satış",
     "alt": "Plasiyerin günde yüzlerce kez açtığı ekran"},

    {"tip": "kartlar",
     "baslik": "Müşteri kapısında altı dokunuş",
     "alt": "Ekran, sahada tek elle kullanılacak şekilde kuruldu",
     "kartlar": [
         ("1 · Müşteriyi seç", "Bakiye, kredi limiti, risk durumu ve son "
          "alışverişleri aynı ekranda görünür."),
         ("2 · Araç stoğunu gör", "Araçta fiilen ne varsa o listelenir; SKT'si "
          "yaklaşan ürün renkli rozetle uyarır."),
         ("3 · Yapay zeka önerisi", "Müşterinin ortalama tüketimi ve son alış "
          "tarihinden hareketle hangi ürünün bitmiş olduğu tahmin edilir."),
         ("4 · Sepeti kur", "Miktar ve birim girilir; kampanya arka planda "
          "kendiliğinden uygulanır, bedava ürün ayrı satır olur."),
         ("5 · Tahsilatı al", "Nakit, kart, havale, çek veya açık hesap. "
          "Çek tahsil edilene kadar bakiyeden düşmez."),
         ("6 · Satışı tamamla", "Sipariş, teslimat, fatura, cari kayıt ve stok "
          "düşümü tek seferde oluşur."),
     ]},

    {"tip": "ekran", "dosya": "03-sicak-satis", "ust": "SICAK SATIŞ",
     "baslik": "Plasiyerin ekranı",
     "maddeler": [
         "Müşteri seçilince bakiye, kredi limiti, risk skoru ve geçmiş alışverişler görünür",
         "Orta sütun aracın gerçek stoğu; SKT'si yaklaşan ürün rozetle uyarır",
         "Kampanya arka planda uygulanır — bedava ürün ayrı satır olarak eklenir",
         "Toplamlar canlı hesaplanır: brüt, satır iskontosu, kampanya, KDV",
         "Ödeme tipi seçilir ve satış tek işlemde tamamlanır",
     ]},

    {"tip": "vurgu",
     "ust": "TASARIM KARARI",
     "baslik": "Ya hepsi yazılır, ya hiçbiri.",
     "alt": "Sipariş → teslimat → FEFO stok çıkışı → fatura → cari kayıt → tahsilat "
            "zinciri tek veritabanı işlemidir. Bir adım başarısız olursa hiçbiri "
            "yazılmaz. Yarım kalmış satış diye bir durum oluşmaz.",
     "kod": "sipariş + teslimat + fatura + cari + tahsilat  =  1 işlem"},

    {"tip": "kartlar",
     "baslik": "Çevrimdışı çalışma",
     "alt": "Sahada kapsama her zaman yoktur — sistem bunu varsayar",
     "kartlar": [
         ("Cihazda tutulur", "Bağlantı yokken sepet cihazda bekler; plasiyer "
          "satışa devam eder."),
         ("Tekrar gönderim güvenli", "Her sepet benzersiz bir kimlik taşır. Aynı "
          "satış iki kez gönderilse bile ikinci kayıt oluşmaz."),
         ("Referans veri önbellekte", "Ürün, müşteri ve kampanya listeleri "
          "çevrimdışı da açılır."),
         ("Kurulabilir uygulama", "PWA olarak telefona/tablete kurulur; ayrı "
          "mağaza süreci gerekmez."),
     ]},

    {"tip": "bolum", "no": "02", "baslik": "Stok",
     "alt": "Gıda dağıtımında stok, tarih demektir"},

    {"tip": "kartlar",
     "baslik": "Değiştirilemez stok defteri",
     "alt": "Gün sonu tartışmasının hakemi olan yapı",
     "kartlar": [
         ("Hareket silinmez", "Stok hareketleri yalnızca eklenir; düzeltme de "
          "yeni bir hareket olarak yazılır. Geçmiş değiştirilemez."),
         ("Bakiye defterden türer", "Hızlı okuma için bakiye ayrıca tutulur ama "
          "her an defterden yeniden üretilebilir."),
         ("FEFO varsayılan", "Son kullanma tarihi en yakın parti önce çıkar. "
          "Gıdada bu bir tercih değil, zorunluluktur."),
         ("Parti izlenebilirliği", "Hangi partinin hangi müşteriye gittiği "
          "kayıtlıdır — geri çağırma anında listelenir."),
     ]},

    {"tip": "ekran", "dosya": "08-arac-yukleme", "ust": "ARAÇ YÜKLEME",
     "baslik": "Yapay zekâ destekli yükleme",
     "maddeler": [
         "Rotadaki müşterilerin geçmiş tüketimi, haftanın günü ve mevsim etkisi",
         "Araçta hâlihazırda bulunan stok düşülür",
         "Hacim ve ağırlık kapasitesi aşılamaz",
         "Her satır kendi gerekçesini ve güven düzeyini taşır",
         "Öneri düzenlenebilir — son söz plasiyerindir",
     ]},

    {"tip": "vurgu",
     "ust": "GÜN SONU MUTABAKATI",
     "baslik": "Fark artık tartışılmaz, hesaplanır.",
     "alt": "Her terim hareket defterinden gelir; elle tutulan hiçbir sayaç yoktur. "
            "Fark sıfırdan farklıysa gün oturumu işaretlenir, düzeltme hareketi "
            "yazılır ve yöneticiye bildirim gider.",
     "kod": "teorik = açılış + yüklenen + ek yükleme − satılan + iade − fire\n"
            "fark   = teorik − fiziki sayım"},

    {"tip": "kartlar",
     "baslik": "Araç = mobil depo",
     "alt": "Ayrı bir 'araç stoğu' alt sistemi yazmadık",
     "kartlar": [
         ("Aynı kurallar", "Araç, tipi VEHICLE olan bir depodur. FEFO, sayım, "
          "transfer ve değerleme araçta da aynen çalışır."),
         ("Otomatik kurulum", "Araç kaydedilince deposu kendiliğinden oluşur; "
          "ayrıca tanımlamak gerekmez."),
         ("Kapasite denetimi", "Yükleme sırasında hacim ve ağırlık kapasitesi "
          "kontrol edilir, aşım engellenir."),
         ("Neden böyle", "İkinci bir stok alt sistemi, aynı kuralların zamanla "
          "ayrışan bir kopyasını üretirdi."),
     ]},

    {"tip": "bolum", "no": "03", "baslik": "Saha",
     "alt": "Rota, ziyaret ve konum"},

    {"tip": "kartlar",
     "baslik": "Rota optimizasyonu",
     "alt": "Kapasite, zaman penceresi ve servis süresi birlikte çözülür",
     "kartlar": [
         ("İki katmanlı çözücü", "Google OR-Tools kuruluysa kesin çözücü; "
          "değilse sistemin kendi Clarke-Wright + 2-opt çözücüsü devreye girer."),
         ("Özellik kaybolmaz", "100 MB'lık bir bağımlılığın kurulu olmaması "
          "rotayı devre dışı bırakmaz — yedek çözücü her zaman çalışır."),
         ("Gerçek kısıtlar", "Araç hacmi ve ağırlığı, müşteri çalışma saatleri, "
          "servis süresi, mesai limiti ve öncelikli müşteriler."),
         ("Planlanan / gerçekleşen", "Tamamlanan, atlanan ve geciken duraklar; "
          "planlanan ve gerçekleşen kilometre farkı raporlanır."),
     ]},

    {"tip": "ekran", "dosya": "10-rotalar", "ust": "SAHA",
     "baslik": "Rota ve harita",
     "maddeler": [
         "Günlük rotalar müşterilerin ziyaret günlerinden üretilir",
         "Optimizasyon kapasite, zaman penceresi ve servis süresini birlikte çözer",
         "Hangi çözücünün çalıştığı ekranda yazar",
         "Planlanan ve gerçekleşen kilometre farkı raporlanır",
     ]},

    {"tip": "bolum", "no": "04", "baslik": "Analitik",
     "alt": "Ölçülemeyen tahmin, güvenilmez tahmindir"},

    {"tip": "kartlar",
     "baslik": "Talep tahmini",
     "alt": "FMCG talebi kesintilidir — klasik yöntemler bu seride kötü çalışır",
     "kartlar": [
         ("Önce sınıflandır", "Seri düzgün mü, kesintili mi, düzensiz mi? "
          "Yöntem bu sınıfa göre seçilir, tek yöntem dayatılmaz."),
         ("Doğru yöntem", "Croston, SBA, TSB, Holt-Winters veya haftanın günü "
          "mevsimsel naif — hangisi bu seriye uyuyorsa."),
         ("Geriye dönük test", "Seçim MAE / MAPE / RMSE ile doğrulanır. "
          "Kullanılan yöntem ve hata payı ekranda gösterilir."),
         ("Sıfırlar sayılır", "Takvim sıfırlarla doldurulur; kesintili talebi "
          "tanımlayan şey zaten o sıfırlardır."),
     ]},

    {"tip": "kartlar",
     "baslik": "Kontrol paneli ve raporlar",
     "alt": "Panel ile rapor aynı toplulaştırmadan beslenir — çelişemezler",
     "kartlar": [
         ("20 KPI", "Günlük/aylık satış, hedef gerçekleşme, brüt kâr, tahsilat, "
          "açık ve geciken alacak, stok değeri, kritik stok, SKT yaklaşan."),
         ("7 canlı grafik", "Satış trendi, kategori kırılımı, en iyi ürün ve "
          "plasiyerler, bölge dağılımı, tahsilat/satış karşılaştırması."),
         ("21 hazır rapor", "Satış, plasiyer, müşteri, SKU, marka, bölge, rota, "
          "tahsilat, risk, stok, SKT, fire, iade, kampanya, kârlılık, hedef."),
         ("PDF · Excel · CSV", "Türkçe karakterler için UTF-8 BOM ile üretilir; "
          "Excel'de içe aktarma adımı gerekmeden düzgün açılır."),
     ]},

    {"tip": "ekran", "dosya": "18-tahminler", "ust": "ANALİTİK",
     "baslik": "Talep tahmini",
     "maddeler": [
         "Seri önce sınıflandırılır: düzgün, kesintili veya düzensiz",
         "Yöntem sınıfa göre seçilir — tek yöntem dayatılmaz",
         "Seçim geriye dönük testle doğrulanır (MAE / MAPE / RMSE)",
         "Kullanılan yöntem, hata payı ve güven aralığı ekranda görünür",
     ]},

    {"tip": "bolum", "no": "05", "baslik": "Yapay Zeka",
     "alt": "Yerel öncelikli, açıklanabilir, bütçeli"},

    {"tip": "kartlar",
     "baslik": "Üç sağlayıcı, tek arayüz",
     "alt": "Yönlendirici göreve göre model seçer ve sağlayıcı düşerse yedeğe geçer",
     "kartlar": [
         ("LM Studio — yerel", "Bilgisayarınızda çalışır, ücret yok, veri dışarı "
          "çıkmaz. Varsayılan ilk tercih."),
         ("NVIDIA NIM — bulut", "Yerel model yetmediğinde devreye girer; "
          "100'den fazla model arasından göreve göre seçim."),
         ("Claude — bulut", "Uzun bağlam ve karmaşık muhakeme gerektiren işler "
          "için yapılandırılabilir."),
         ("Bütçe koruması", "Aylık limit aşılınca ücretli sağlayıcılar durur, "
          "yerel model çalışmaya devam eder."),
     ]},

    {"tip": "kartlar",
     "baslik": "AI Satış Müdürü",
     "alt": "Doğal Türkçe soru → salt okunur sorgu → gerekçeli cevap",
     "kartlar": [
         ("Örnek sorular", "«Bugün en fazla satış yapan 10 plasiyeri getir» · "
          "«Tahsilat riski yüksek müşterileri göster» · «Son 90 günde "
          "kaybettiğimiz müşteriler»"),
         ("Kaynağı gösterir", "Cevabın hangi veriden üretildiği açılır panelde "
          "görünür — rakam nereden geldi sorusu cevapsız kalmaz."),
         ("Yalnızca okuma", "Tek bir SELECT çalıştırılabilir. Veri değiştiren "
          "hiçbir komut kabul edilmez."),
         ("Kapalı tablolar", "Kullanıcı, oturum ve denetim tabloları yapay zeka "
          "sorgularına tamamen kapalıdır."),
     ]},

    {"tip": "ekran", "dosya": "19-ai-mudur", "ust": "YAPAY ZEKA",
     "baslik": "AI Satış Müdürü",
     "maddeler": [
         "Doğal Türkçe soru sorulur",
         "Sistem yalnızca okuma sorgusu üretir ve çalıştırır",
         "Cevabın hangi veriden üretildiği açılır panelde görünür",
         "Kullanıcı, oturum ve denetim tabloları erişime kapalıdır",
     ]},

    {"tip": "kartlar",
     "baslik": "Güvenlik",
     "alt": "Yetki üç katmanda ayrı ayrı denetlenir",
     "kartlar": [
         ("%s rol · %s izin" % (_s("rol"), _s("izin")), "Ekran, işlem ve veri kapsamı ayrı ayrı "
          "denetlenir. Plasiyer yalnızca kendi kayıtlarını görür."),
         ("Yetki yükseltme kapalı", "Kimse kendinden üstün bir rol atayamaz, "
          "sahip olmadığı bir yetkiyi başkasına veremez."),
         ("Zincirli denetim kaydı", "Her kayıt bir öncekinin özetini taşır. "
          "Geçmiş bir satır değişirse zincir kırılır ve yeri bildirilir."),
         ("Sır sızmaz", "API anahtarları veritabanında saklanmaz, loglara "
          "yazılmadan temizlenir, ekranda maskelenir."),
     ]},

    {"tip": "ekran", "dosya": "22-uyumluluk", "ust": "UYUMLULUK",
     "baslik": "Uyumluluk ve insan hakları katmanı",
     "maddeler": [
         "Kişisel veri envanteri koddan ölçülerek çıkarılır",
         "Rıza ve aydınlatma ayrı kayıtlar olarak tutulur",
         "Kanıt eksikse sonuç 'uyumlu' değil, 'kanıt yetersiz' olur",
         "Yetki motoruna sorulan her karar — reddedilenler dâhil — zincirli bir makbuz ve itiraz yolu bırakır",
     ]},

    {"tip": "galeri",
     "baslik": "Diğer ekranlar",
     "alt": "Sistem %s ekrandan oluşur; burada altı tanesi" % _s("ekran"),
     "ogeler": [
         ("04-musteriler", "Müşteriler"),
         ("06-urunler", "Ürünler"),
         ("16-raporlar", "Raporlar"),
         ("12-gun-yonetimi", "Gün yönetimi"),
         ("28-roller", "Roller ve yetkiler"),
         ("27-egitim", "Eğitim merkezi"),
     ]},

    {"tip": "iki_sutun",
     "baslik": "Teknoloji ve kurulum",
     "alt": "Kurulu olmayan hiçbir şey zorunlu kılınmadı",
     "sol_baslik": "Teknoloji",
     "sol": ["Python 3.11 · FastAPI · SQLAlchemy 2.0",
             "React 18 · TypeScript · Vite · Tailwind",
             "SQLite varsayılan — PostgreSQL opsiyonel",
             "PWA: çevrimdışı çalışma, kurulabilir",
             "Redis ve Docker zorunlu değil"],
     "sag_baslik": "Kurulum",
     "sag": ["setup.ps1 — tek komutla kurulum",
             "start.bat — tek tıkla başlatma",
             "Demo veri: 500 müşteri, 12 aylık geçmiş",
             "14 derslik uygulama içi eğitim merkezi",
             "TR / EN — anında dil değişimi"]},

    {"tip": "vurgu",
     "ust": "PLATFORMLAR & İNDİRME",
     "baslik": "Windows 10/11 (x64) + macOS (Apple Silicon)",
     "alt": "GitHub Releases üzerinde iki paket yayımlandı: Windows zip + macOS zip. "
            "macOS paketi Apple Silicon (arm64) içindir ve notarize edilmemiştir — "
            "ilk açılışta sağ tık → Aç yeterlidir.",
     "kod": "İndirme: github.com/Azizsekerdil/smart-van-sales-management-system/releases\n"
            "Sürüm: v1.0.0 — Windows zip · macOS zip"},

    {"tip": "kapanis",
     "baslik": "Akıllı Sıcak Satış Yönetim Sistemi",
     "alt": "Sahadaki her hareketi kayda geçiren, gün sonunda hesabı veren sistem",
     "maddeler": ["Kurulum: setup.ps1", "Başlatma: start.bat",
                  "Kılavuz: docs/KULLANIM_KILAVUZU_TR.md"],
     "dip": "Bu sunum programın kendi sürümünden üretilmiştir"},
]


SLAYTLAR_EN = [
    {"tip": "kapak",
     "baslik": "Smart Van Sales\nManagement System",
     "alt": "End-to-end direct-store-delivery platform for food & beverage",
     "etiketler": ["Hot Sale", "FEFO Stock", "Route Optimisation", "Local-first AI"],
     "dip": "From the factory to the shop shelf — one system"},

    {"tip": "vurgu",
     "ust": "THE PROBLEM",
     "baslik": "At day end, everyone quotes\na different number.",
     "alt": "When van stock, invoices raised, cash collected and the physical count "
            "are each kept separately, nobody can trace where the evening's "
            "discrepancy came from."},

    {"tip": "kartlar",
     "baslik": "Four places it breaks",
     "alt": "Where van sales operations most reliably lose money",
     "kartlar": [
         ("Guessed van stock", "What was loaded is in a notebook, what was sold is "
          "on receipts. When they disagree at night, the variance belongs to nobody."),
         ("Expiry leakage", "Short-dated stock sits on the shelf while a newer batch "
          "is sold. The write-off is discovered only at disposal."),
         ("Route blindness", "Routes are built from habit. Nobody measures the "
          "relationship between kilometres driven and revenue earned."),
         ("Credit risk", "The credit limit lives on paper; the office learns about "
          "the over-limit sale after the rep has already made it."),
     ]},

    {"tip": "zincir",
     "baslik": "The chain this system covers",
     "alt": "Every link feeds from the previous one — no manual re-entry anywhere",
     "adimlar": ["Factory", "Central Depot", "Regional Depot", "Van Loading",
                 "Sales Van", "Salesperson", "Daily Route", "Customer",
                 "Hot Sale", "Invoice / Waybill", "Collection", "Return",
                 "Van Count", "Day-End Reconciliation"]},

    {"tip": "sayilar",
     "baslik": "The measure of the system",
     "alt": "Figures measured from the running system, not targets",
     "sayilar": [(_s("tablo"), "database tables"), (_s("uc_nokta"), "API endpoints"),
                 (_s("ekran"), "interface screens"), (_s("izin"), "permissions"),
                 (_s("servis"), "business services"), (_s("rapor"), "built-in reports"),
                 (_s("indeks"), "database indexes"), (_s("dil"), "languages (TR / EN)")]},

    {"tip": "ekran_tam", "dosya": "02-panel",
     "baslik": "The dashboard",
     "alt": "20 KPIs and 7 live charts — panel and reports share one aggregation"},

    {"tip": "bolum", "no": "01", "baslik": "Hot Sale",
     "alt": "The screen a rep opens hundreds of times a day"},

    {"tip": "kartlar",
     "baslik": "Six taps at the customer's door",
     "alt": "Built to be used one-handed, standing in a shop",
     "kartlar": [
         ("1 · Pick the customer", "Balance, credit limit, risk and recent "
          "purchases all on the same screen."),
         ("2 · See the van stock", "Only what is physically on the van; "
          "short-dated items carry a coloured warning badge."),
         ("3 · Read the AI suggestion", "From average consumption and days since "
          "the last purchase, the system estimates what has run out."),
         ("4 · Build the basket", "Enter quantity and unit; campaigns apply "
          "themselves and free goods appear as their own line."),
         ("5 · Take payment", "Cash, card, transfer, cheque or open account. "
          "A cheque does not reduce the balance until it clears."),
         ("6 · Complete the sale", "Order, delivery, invoice, ledger entry and "
          "stock issue are created in one go."),
     ]},

    {"tip": "ekran", "dosya": "03-sicak-satis", "ust": "HOT SALE",
     "baslik": "The salesperson's screen",
     "maddeler": [
         "Picking a customer shows balance, credit limit, risk score and past orders",
         "The middle column is the van's real stock; short-dated items carry a badge",
         "Campaigns apply themselves — free goods appear as their own line",
         "Totals compute live: gross, line discount, campaign, VAT",
         "Payment method is chosen and the sale posts in one transaction",
     ]},

    {"tip": "vurgu",
     "ust": "DESIGN DECISION",
     "baslik": "All of it is written, or none of it.",
     "alt": "Order → delivery → FEFO stock issue → invoice → ledger → payment is a "
            "single database transaction. If any step fails, none of it is written. "
            "There is no such thing as a half-finished sale.",
     "kod": "order + delivery + invoice + ledger + payment  =  1 transaction"},

    {"tip": "kartlar",
     "baslik": "Working offline",
     "alt": "Coverage is not guaranteed in the field — the system assumes that",
     "kartlar": [
         ("Held on the device", "With no connection the basket waits on the "
          "device and the rep keeps selling."),
         ("Safe to resend", "Every basket carries a unique id. Submitting the "
          "same sale twice never creates a second record."),
         ("Reference data cached", "Product, customer and campaign lists open "
          "without a connection."),
         ("Installable app", "Installs to a phone or tablet as a PWA — no app "
          "store process required."),
     ]},

    {"tip": "bolum", "no": "02", "baslik": "Stock",
     "alt": "In food distribution, stock means dates"},

    {"tip": "kartlar",
     "baslik": "An immutable stock ledger",
     "alt": "The structure that settles the day-end argument",
     "kartlar": [
         ("Movements are never deleted", "Stock movements are only appended; a "
          "correction is itself a new movement. History cannot be rewritten."),
         ("Balance derives from the ledger", "A balance is kept for fast reads "
          "but can be rebuilt from the ledger at any moment."),
         ("FEFO by default", "The nearest-expiry batch leaves first. In food "
          "that is not a preference, it is a requirement."),
         ("Batch traceability", "Which batch went to which customer is recorded "
          "— a recall produces a list immediately."),
     ]},

    {"tip": "ekran", "dosya": "08-arac-yukleme", "ust": "VAN LOADING",
     "baslik": "AI-assisted loading",
     "maddeler": [
         "Past consumption of the route's customers, weekday and season effects",
         "Stock already on the van is subtracted",
         "Volume and weight capacity cannot be exceeded",
         "Every line carries its own reason and confidence level",
         "The suggestion is editable — the rep has the last word",
     ]},

    {"tip": "vurgu",
     "ust": "DAY-END RECONCILIATION",
     "baslik": "The variance is computed, not debated.",
     "alt": "Every term comes from the movement ledger; there is no hand-maintained "
            "counter anywhere. A non-zero variance flags the session, writes an "
            "adjustment movement and notifies the manager.",
     "kod": "theoretical = opening + loaded + reloaded − sold + returned − wastage\n"
            "variance    = theoretical − physical count"},

    {"tip": "kartlar",
     "baslik": "A van is a warehouse",
     "alt": "We did not write a separate 'van stock' subsystem",
     "kartlar": [
         ("Same rules", "A van is a warehouse of type VEHICLE. FEFO, counting, "
          "transfers and valuation all work on it unchanged."),
         ("Created automatically", "Saving a vehicle creates its warehouse; "
          "there is nothing extra to set up."),
         ("Capacity enforced", "Volume and weight are checked during loading and "
          "an overload is refused."),
         ("Why this way", "A second stock subsystem would be a copy of the same "
          "rules that drifts apart over time."),
     ]},

    {"tip": "bolum", "no": "03", "baslik": "The Field",
     "alt": "Routes, visits and position"},

    {"tip": "kartlar",
     "baslik": "Route optimisation",
     "alt": "Capacity, time windows and service time solved together",
     "kartlar": [
         ("Two-layer solver", "Google OR-Tools when installed; otherwise the "
          "system's own Clarke-Wright + 2-opt solver takes over."),
         ("The feature never disappears", "A missing 100 MB dependency does not "
          "disable routing — the fallback solver always works."),
         ("Real constraints", "Vehicle volume and weight, customer opening hours, "
          "service time, working-day limit and priority customers."),
         ("Plan vs actual", "Completed, skipped and delayed stops, plus the "
          "planned-versus-actual kilometre difference."),
     ]},

    {"tip": "ekran", "dosya": "10-rotalar", "ust": "THE FIELD",
     "baslik": "Routes and map",
     "maddeler": [
         "Daily routes are generated from customers' visit days",
         "Optimisation solves capacity, time windows and service time together",
         "Which solver ran is stated on screen",
         "Planned versus actual kilometres are reported",
     ]},

    {"tip": "bolum", "no": "04", "baslik": "Analytics",
     "alt": "A forecast you cannot measure is a forecast you cannot trust"},

    {"tip": "kartlar",
     "baslik": "Demand forecasting",
     "alt": "FMCG demand is intermittent — classical methods do badly on it",
     "kartlar": [
         ("Classify first", "Is the series smooth, intermittent or lumpy? The "
          "method follows from the class; no single method is imposed."),
         ("Then choose", "Croston, SBA, TSB, Holt-Winters or a day-of-week "
          "seasonal naive — whichever suits this particular series."),
         ("Back-test the choice", "Validated with MAE / MAPE / RMSE. The method "
          "used and its error are shown on screen."),
         ("Zeros count", "The calendar is zero-filled; those zeros are what "
          "define intermittent demand in the first place."),
     ]},

    {"tip": "kartlar",
     "baslik": "Dashboard and reports",
     "alt": "Panel and report share one aggregation — they cannot disagree",
     "kartlar": [
         ("20 KPIs", "Daily and monthly sales, target achievement, gross margin, "
          "collections, open and overdue receivables, stock value, expiring soon."),
         ("7 live charts", "Sales trend, category split, top products and "
          "salespeople, regional distribution, collections versus sales."),
         ("21 built-in reports", "Sales, salesperson, customer, SKU, brand, "
          "region, route, collections, risk, stock, expiry, wastage, returns, "
          "campaigns, profitability, targets."),
         ("PDF · Excel · CSV", "Written as UTF-8 with BOM so Turkish characters "
          "open correctly in Excel with no import step."),
     ]},

    {"tip": "ekran", "dosya": "18-tahminler", "ust": "ANALYTICS",
     "baslik": "Demand forecasting",
     "maddeler": [
         "The series is classified first: smooth, intermittent or lumpy",
         "The method follows from the class — none is imposed",
         "The choice is validated by back-testing (MAE / MAPE / RMSE)",
         "Method used, error and confidence interval are shown on screen",
     ]},

    {"tip": "bolum", "no": "05", "baslik": "Artificial Intelligence",
     "alt": "Local-first, explainable, budgeted"},

    {"tip": "kartlar",
     "baslik": "Three providers, one interface",
     "alt": "The router picks a model per task and fails over when one is down",
     "kartlar": [
         ("LM Studio — local", "Runs on your own machine, costs nothing and no "
          "data leaves the building. The default first choice."),
         ("NVIDIA NIM — cloud", "Steps in when the local model is not enough; "
          "over 100 models selected by task type."),
         ("Claude — cloud", "Configurable for long-context and harder reasoning "
          "work."),
         ("Budget protection", "When the monthly limit is reached the paid "
          "providers stop while the local model keeps working."),
     ]},

    {"tip": "kartlar",
     "baslik": "The AI Sales Manager",
     "alt": "Plain-language question → read-only query → justified answer",
     "kartlar": [
         ("Example questions", "«Show the 10 best-selling reps today» · «Show "
          "customers with high collection risk» · «Which customers did we lose "
          "in the last 90 days?»"),
         ("It shows its source", "The data behind the answer is visible in an "
          "expandable panel — \"where did that number come from\" is answerable."),
         ("Read-only", "A single SELECT may be executed. No data-modifying "
          "command is accepted."),
         ("Closed tables", "User, session and audit tables are entirely off "
          "limits to AI queries."),
     ]},

    {"tip": "ekran", "dosya": "19-ai-mudur", "ust": "ARTIFICIAL INTELLIGENCE",
     "baslik": "The AI Sales Manager",
     "maddeler": [
         "A question is asked in plain language",
         "The system produces and runs a read-only query only",
         "The data behind the answer is visible in an expandable panel",
         "User, session and audit tables are off limits",
     ]},

    {"tip": "kartlar",
     "baslik": "Security",
     "alt": "Authorisation is checked separately at three layers",
     "kartlar": [
         ("%s roles · %s permissions" % (_s("rol"), _s("izin")), "Screen, action and data scope are "
          "enforced separately. A rep sees only their own records."),
         ("No privilege escalation", "Nobody can assign a role above their own "
          "or grant a permission they do not hold."),
         ("Chained audit log", "Each entry carries the previous one's digest. "
          "Alter a historical row and the chain breaks, reporting where."),
         ("Secrets stay secret", "API keys are never stored in the database, are "
          "stripped before logging and are masked on screen."),
     ]},

    {"tip": "ekran", "dosya": "22-uyumluluk", "ust": "COMPLIANCE",
     "baslik": "Compliance and human rights layer",
     "maddeler": [
         "The personal-data inventory is measured from the code, not assumed",
         "Consent and notice are kept as separate records",
         "Missing evidence yields 'insufficient evidence', never 'compliant'",
         "Every decision put to the authority engine — refusals included — leaves a chained receipt and an appeal path",
     ]},

    {"tip": "galeri",
     "baslik": "Other screens",
     "alt": "The system has %s screens; six of them here" % _s("ekran"),
     "ogeler": [
         ("04-musteriler", "Customers"),
         ("06-urunler", "Products"),
         ("16-raporlar", "Reports"),
         ("12-gun-yonetimi", "Day sessions"),
         ("28-roller", "Roles & permissions"),
         ("27-egitim", "Training centre"),
     ]},

    {"tip": "iki_sutun",
     "baslik": "Technology and setup",
     "alt": "Nothing that was not already installed has been made mandatory",
     "sol_baslik": "Technology",
     "sol": ["Python 3.11 · FastAPI · SQLAlchemy 2.0",
             "React 18 · TypeScript · Vite · Tailwind",
             "SQLite by default — PostgreSQL optional",
             "PWA: offline capable, installable",
             "Redis and Docker are not required"],
     "sag_baslik": "Setup",
     "sag": ["setup.ps1 — one-command install",
             "start.bat — one-click start",
             "Demo data: 500 customers, 12 months history",
             "14-lesson in-app training centre",
             "TR / EN — instant language switch"]},

    {"tip": "vurgu",
     "ust": "PLATFORMS & DOWNLOAD",
     "baslik": "Windows 10/11 (x64) & macOS (Apple Silicon)",
     "alt": "Two assets are published on GitHub Releases: a Windows zip and a macOS zip. "
            "The macOS package targets Apple Silicon (arm64) and is not notarized — "
            "on first launch, right-click → Open.",
     "kod": "Download: GitHub Releases (v1.0.0)\n"
            "github.com/Azizsekerdil/smart-van-sales-management-system/releases"},

    {"tip": "kapanis",
     "baslik": "Smart Van Sales Management System",
     "alt": "Every movement in the field recorded, every figure accounted for at day end",
     "maddeler": ["Install: setup.ps1", "Start: start.bat",
                  "Guide: docs/USER_GUIDE_EN.md"],
     "dip": "This deck is generated from the product's own version"},
]

SLAYTLAR = {"tr": SLAYTLAR_TR, "en": SLAYTLAR_EN}

METIN = {
    "tr": {"marka": "Akıllı Sıcak Satış Yönetim Sistemi", "slayt": "Slayt",
           "baslik": "Akıllı Sıcak Satış Yönetim Sistemi — Tanıtım",
           "onceki": "Önceki", "sonraki": "Sonraki", "tamekran": "Tam ekran",
           "kapat": "Kapat", "kod": "tr",
           "ipucu": "Kaydırarak ilerleyin · ok tuşları da çalışır",
           "uyari": "Slaytlar yüklenemedi."},
    "en": {"marka": "Smart Van Sales Management System", "slayt": "Slide",
           "baslik": "Smart Van Sales Management System — Introduction",
           "onceki": "Previous", "sonraki": "Next", "tamekran": "Full screen",
           "kapat": "Close", "kod": "en",
           "ipucu": "Swipe to advance · arrow keys work too",
           "uyari": "Slides could not be loaded."},
}


# ══════════════════════════════════════════════════════════════════════════
# YOLLAR
# ══════════════════════════════════════════════════════════════════════════
#: Yayın sürümü eki. Üretilen her dosya bunu taşır, böylece bir dosyanın
#: denetlenmiş sürüm olup olmadığı adından anlaşılır.
KAMUYA_ACIK_EK = "_PUBLIC"


def _ad(dil):
    taban = "Van_Sales_Tanitim" if dil == "tr" else "Van_Sales_Intro_EN"
    return taban + KAMUYA_ACIK_EK


def _yollar(dil):
    t = os.path.join(CIKTI, _ad(dil))
    return t + ".pptx", t + ".pdf", t + ".html"


def _baski_yollari(dil):
    ek = "_Baski" if dil == "tr" else "_Print"
    t = os.path.join(CIKTI, _ad(dil) + ek)
    return t + ".pptx", t + ".pdf"


# ══════════════════════════════════════════════════════════════════════════
# ÇİZİM YARDIMCILARI
# ══════════════════════════════════════════════════════════════════════════
def _renk(hexstr):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hexstr)


def _kutu(slayt, x, y, w, h, dolgu=None, cizgi=None, yuvarlak=True, kalinlik=1.0):
    """Gölgesiz dikdörtgen/yuvarlak kutu. Gölge kapatılır: 20+ şekilli slaytta
    PowerPoint'in varsayılan gölgesi kirli bir katman oluşturuyor."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    sekil = slayt.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if yuvarlak else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sekil.shadow.inherit = False
    if yuvarlak:
        try:
            sekil.adjustments[0] = 0.06
        except Exception:
            pass
    if dolgu:
        sekil.fill.solid()
        sekil.fill.fore_color.rgb = _renk(dolgu)
    else:
        sekil.fill.background()
    if cizgi:
        sekil.line.color.rgb = _renk(cizgi)
        sekil.line.width = Pt(kalinlik)
    else:
        sekil.line.fill.background()
    sekil.text_frame.text = ""
    return sekil


def _yazi(slayt, x, y, w, h, satirlar, hiza="l", dikey="t"):
    """satirlar: [(metin, punto, renk, kalın, satır_aralığı_ekstra_pt)]"""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    kutu = slayt.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = kutu.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[dikey]

    for i, (metin, punto, renk, kalin, bosluk) in enumerate(satirlar):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[hiza]
        if bosluk:
            p.space_before = Pt(bosluk)
        # \n içeren metinler ayrı satır olarak kurulur; python-pptx tek run içinde
        # satır sonu yönetmediği için parçalayıp <a:br> yerine paragraf kullanıyoruz.
        for j, parca in enumerate(str(metin).split("\n")):
            if j > 0:
                p = tf.add_paragraph()
                p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                               "r": PP_ALIGN.RIGHT}[hiza]
            r = p.add_run()
            r.text = parca
            r.font.size = Pt(punto)
            r.font.bold = kalin
            r.font.color.rgb = _renk(renk)
            r.font.name = YAZI_TIPI
    return kutu


def _zemin(slayt, renk=None):
    slayt.background.fill.solid()
    slayt.background.fill.fore_color.rgb = _renk(renk or Z["zemin"])


def _ustbilgi(slayt, marka, surum, sayfa, toplam):
    """Her içerik slaydının üst şeridi — marka solda, sayfa numarası sağda."""
    _yazi(slayt, 0.62, 0.34, 8.0, 0.3,
          [(marka.upper(), 10, Z["sonuk2"], True, 0)])
    _yazi(slayt, 10.2, 0.34, 2.5, 0.3,
          [("%s   %02d / %02d" % (surum, sayfa, toplam), 10, Z["sonuk2"], False, 0)],
          hiza="r")
    _kutu(slayt, 0.62, 0.68, 12.1, 0.014, dolgu=Z["kart_cizgi"], yuvarlak=False)


def _slayt_basligi(slayt, baslik, alt):
    _yazi(slayt, 0.62, 1.02, 11.4, 0.62, [(baslik, 30, Z["yazi_ac"], True, 0)])
    if alt:
        _yazi(slayt, 0.62, 1.72, 11.4, 0.4, [(alt, 13, Z["sonuk"], False, 0)])


# ══════════════════════════════════════════════════════════════════════════
# SLAYT TİPLERİ
# ══════════════════════════════════════════════════════════════════════════
def _s_kapak(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    # Sol kenarda marka şeridi — kapağı içerik slaytlarından ayıran tek öğe.
    _kutu(slayt, 0, 0, 0.22, 7.5, dolgu=Z["marka"], yuvarlak=False)
    _kutu(slayt, 7.9, -1.2, 7.0, 7.0, dolgu=Z["zemin2"], yuvarlak=True)

    _yazi(slayt, 1.05, 1.30, 7.4, 0.4, [(marka.upper(), 11, Z["vurgu"], True, 0)])
    _yazi(slayt, 1.05, 1.85, 8.2, 2.1, [(d["baslik"], 47, Z["yazi_ac"], True, 0)])
    _yazi(slayt, 1.05, 4.05, 7.6, 0.8, [(d["alt"], 15, Z["sonuk"], False, 0)])

    x = 1.05
    for etiket in d["etiketler"]:
        w = 0.30 + len(etiket) * 0.098
        _kutu(slayt, x, 5.05, w, 0.42, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _yazi(slayt, x, 5.05, w, 0.42, [(etiket, 11, Z["marka_ac"], True, 0)],
              hiza="c", dikey="m")
        x += w + 0.16

    _kutu(slayt, 1.05, 6.10, 0.9, 0.045, dolgu=Z["vurgu"], yuvarlak=False)
    _yazi(slayt, 1.05, 6.35, 8.0, 0.4, [(d["dip"], 12, Z["sonuk2"], False, 0)])
    _yazi(slayt, 10.0, 6.35, 2.7, 0.4, [(surum, 12, Z["vurgu"], True, 0)], hiza="r")


def _s_bolum(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _kutu(slayt, 0, 0, 0.22, 7.5, dolgu=Z["marka"], yuvarlak=False)
    _yazi(slayt, 1.15, 2.45, 3.0, 1.5, [(d["no"], 96, Z["kart"], True, 0)])
    _yazi(slayt, 3.30, 2.72, 8.6, 1.0, [(d["baslik"], 42, Z["yazi_ac"], True, 0)])
    _kutu(slayt, 3.34, 3.86, 1.1, 0.045, dolgu=Z["vurgu"], yuvarlak=False)
    _yazi(slayt, 3.30, 4.12, 8.4, 0.6, [(d["alt"], 15, Z["sonuk"], False, 0)])
    _yazi(slayt, 10.2, 6.72, 2.5, 0.3,
          [("%s   %02d / %02d" % (surum, sayfa, toplam), 10, Z["sonuk2"], False, 0)],
          hiza="r")


def _s_vurgu(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _yazi(slayt, 0.62, 1.30, 11.4, 0.35, [(d["ust"], 11, Z["vurgu"], True, 0)])
    _yazi(slayt, 0.62, 1.80, 11.4, 1.5, [(d["baslik"], 36, Z["yazi_ac"], True, 0)])

    kod = d.get("kod")
    alt_y = 3.55
    if kod:
        satir = kod.count("\n") + 1
        h = 0.52 + satir * 0.36
        _kutu(slayt, 0.62, alt_y, 12.1, h, dolgu=Z["kart"], cizgi=Z["marka"])
        _yazi(slayt, 0.98, alt_y + 0.26, 11.4, h - 0.5,
              [(kod, 16, Z["vurgu"], True, 0)])
        alt_y += h + 0.42
    _yazi(slayt, 0.62, alt_y, 11.0, 1.5, [(d["alt"], 14, Z["sonuk"], False, 0)])


def _s_kartlar(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _slayt_basligi(slayt, d["baslik"], d.get("alt"))

    kartlar = d["kartlar"]
    sutun = 2 if len(kartlar) <= 4 else 3
    satir = (len(kartlar) + sutun - 1) // sutun
    bosluk = 0.30
    w = (12.1 - bosluk * (sutun - 1)) / sutun

    # Kart yuksekligi icerikten turetilir. Sabit yukseklik kisa metinlerde
    # kartin altinda genis bir olu bosluk birakiyordu; en uzun govde metnine
    # gore satir sayisi kestirilip ona gore olculuyor.
    kar_sat = max(len(g) / (26.0 * (w - 0.62)) for _, g in kartlar)
    h = max(1.35, min(2.05, 0.98 + kar_sat * 0.20))

    # Blok, baslik ile alt kenar arasindaki bosluga dikeyde ortalanir.
    alan_ust, alan_alt = 2.34, 7.02
    toplam_h = satir * h + (satir - 1) * bosluk
    ust = alan_ust + max(0.0, (alan_alt - alan_ust - toplam_h) / 2)

    for i, (bas, gov) in enumerate(kartlar):
        cx = 0.62 + (i % sutun) * (w + bosluk)
        cy = ust + (i // sutun) * (h + bosluk)
        _kutu(slayt, cx, cy, w, h, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _kutu(slayt, cx, cy, 0.055, h, dolgu=Z["marka"], yuvarlak=False)
        _yazi(slayt, cx + 0.34, cy + 0.24, w - 0.62, 0.42,
              [(bas, 14 if sutun == 2 else 13, Z["yazi_ac"], True, 0)])
        _yazi(slayt, cx + 0.34, cy + 0.72, w - 0.62, h - 0.92,
              [(gov, 11 if sutun == 2 else 10, Z["sonuk"], False, 0)])


def _s_sayilar(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _slayt_basligi(slayt, d["baslik"], d.get("alt"))

    sayilar = d["sayilar"]
    sutun = 4
    satir = (len(sayilar) + sutun - 1) // sutun
    bosluk = 0.30
    w = (12.1 - bosluk * (sutun - 1)) / sutun
    h = 1.62
    ust = 2.65

    for i, (deger, etiket) in enumerate(sayilar):
        cx = 0.62 + (i % sutun) * (w + bosluk)
        cy = ust + (i // sutun) * (h + bosluk)
        _kutu(slayt, cx, cy, w, h, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _yazi(slayt, cx, cy + 0.26, w, 0.75, [(deger, 40, Z["vurgu"], True, 0)],
              hiza="c")
        _yazi(slayt, cx + 0.18, cy + 1.06, w - 0.36, 0.45,
              [(etiket, 11, Z["sonuk"], False, 0)], hiza="c")


def _s_zincir(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _slayt_basligi(slayt, d["baslik"], d.get("alt"))

    adimlar = d["adimlar"]
    sutun = 7
    satir = (len(adimlar) + sutun - 1) // sutun
    bosluk = 0.22
    w = (12.1 - bosluk * (sutun - 1)) / sutun
    h = 1.05
    satir_bosluk = 0.78
    toplam_h = satir * h + (satir - 1) * satir_bosluk
    ust = 2.34 + max(0.0, (7.02 - 2.34 - toplam_h) / 2)

    for i, ad in enumerate(adimlar):
        cx = 0.62 + (i % sutun) * (w + bosluk)
        cy = ust + (i // sutun) * (h + satir_bosluk)
        # İlk ve son halka vurgulanır: zincirin nerede başlayıp bittiği
        # bir bakışta okunsun.
        ilk_son = i == 0 or i == len(adimlar) - 1
        _kutu(slayt, cx, cy, w, h,
              dolgu=Z["marka"] if ilk_son else Z["kart"],
              cizgi=Z["marka"] if ilk_son else Z["kart_cizgi"])
        _yazi(slayt, cx + 0.10, cy, w - 0.20, h,
              [(ad, 10, Z["yazi_ac"] if ilk_son else Z["yazi"], ilk_son, 0)],
              hiza="c", dikey="m")
        if (i % sutun) != sutun - 1 and i != len(adimlar) - 1:
            _yazi(slayt, cx + w, cy, bosluk, h,
                  [("›", 15, Z["sonuk2"], True, 0)], hiza="c", dikey="m")


def _s_iki_sutun(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _slayt_basligi(slayt, d["baslik"], d.get("alt"))

    w = (12.1 - 0.34) / 2
    for k, (bas, maddeler, x) in enumerate((
            (d["sol_baslik"], d["sol"], 0.62),
            (d["sag_baslik"], d["sag"], 0.62 + w + 0.34))):
        _kutu(slayt, x, 2.42, w, 3.95, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _kutu(slayt, x, 2.42, w, 0.055,
              dolgu=Z["marka"] if k == 0 else Z["vurgu"], yuvarlak=False)
        _yazi(slayt, x + 0.42, 2.72, w - 0.84, 0.45, [(bas, 16, Z["yazi_ac"], True, 0)])
        satirlar = []
        for i, m in enumerate(maddeler):
            satirlar.append(("·  " + m, 12, Z["sonuk"], False, 0 if i == 0 else 13))
        _yazi(slayt, x + 0.42, 3.32, w - 0.84, 2.9, satirlar)


def _s_kapanis(prs, slayt, d, surum, marka, sayfa, toplam):
    _zemin(slayt)
    _kutu(slayt, 0, 0, 0.22, 7.5, dolgu=Z["marka"], yuvarlak=False)
    _yazi(slayt, 1.05, 2.30, 11.0, 1.0, [(d["baslik"], 34, Z["yazi_ac"], True, 0)])
    _kutu(slayt, 1.09, 3.42, 1.1, 0.045, dolgu=Z["vurgu"], yuvarlak=False)
    _yazi(slayt, 1.05, 3.70, 10.5, 0.6, [(d["alt"], 15, Z["sonuk"], False, 0)])

    x = 1.05
    for m in d["maddeler"]:
        w = 0.44 + len(m) * 0.092
        _kutu(slayt, x, 4.62, w, 0.52, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _yazi(slayt, x, 4.62, w, 0.52, [(m, 12, Z["yazi"], False, 0)],
              hiza="c", dikey="m")
        x += w + 0.18

    _yazi(slayt, 1.05, 6.35, 8.5, 0.4, [(d["dip"], 11, Z["sonuk2"], False, 0)])
    _yazi(slayt, 10.0, 6.35, 2.7, 0.4, [(surum, 12, Z["vurgu"], True, 0)], hiza="r")



def _ekran_yolu(dosya, dil):
    """Dile göre ekran görüntüsü yolu; yoksa None döner."""
    yol = os.path.join(EKRANLAR, "%s-%s.png" % (dosya, dil))
    return yol if os.path.exists(yol) else None


def _gorsel(slayt, yol, x, y, w):
    """
    Görüntüyü 16:9 oranını koruyarak yerleştirir ve ince bir çerçeve çizer.

    Çerçeve gerekli: koyu slayt zemininde açık renkli arayüz görüntüsü
    kenarsız bırakılınca "yüzüyor" gibi duruyor.
    """
    from pptx.util import Inches

    h = w * 9.0 / 16.0
    resim = slayt.shapes.add_picture(yol, Inches(x), Inches(y), width=Inches(w))
    cerceve = _kutu(slayt, x - 0.035, y - 0.035, w + 0.07, h + 0.07,
                    dolgu=None, cizgi=Z["kart_cizgi"], yuvarlak=False, kalinlik=1.25)
    # Çerçeve resmin ARKASINDA kalmalı, yoksa görüntüyü kapatır.
    slayt.shapes._spTree.remove(cerceve._element)
    slayt.shapes._spTree.insert(list(slayt.shapes._spTree).index(resim._element), cerceve._element)
    return resim


def _s_ekran(prs, slayt, d, surum, marka, sayfa, toplam):
    """Solda açıklama, sağda ekran görüntüsü."""
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)

    yol = _ekran_yolu(d["dosya"], d.get("_dil", "tr"))
    _yazi(slayt, 0.62, 1.30, 3.9, 0.35, [(d.get("ust", "EKRAN"), 10, Z["vurgu"], True, 0)])
    _yazi(slayt, 0.62, 1.72, 3.9, 1.0, [(d["baslik"], 24, Z["yazi_ac"], True, 0)])
    _kutu(slayt, 0.62, 2.86, 0.7, 0.04, dolgu=Z["marka"], yuvarlak=False)

    satirlar = []
    for i, m in enumerate(d.get("maddeler", [])):
        satirlar.append(("·  " + m, 11, Z["sonuk"], False, 0 if i == 0 else 11))
    if satirlar:
        _yazi(slayt, 0.62, 3.15, 3.9, 3.4, satirlar)

    if yol:
        _gorsel(slayt, yol, 4.95, 1.62, 7.75)
    else:
        _kutu(slayt, 4.95, 1.62, 7.75, 4.36, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _yazi(slayt, 4.95, 1.62, 7.75, 4.36,
              [("Ekran görüntüsü bulunamadı", 12, Z["sonuk2"], False, 0)],
              hiza="c", dikey="m")


def _s_ekran_tam(prs, slayt, d, surum, marka, sayfa, toplam):
    """Üstte başlık, altta tam genişlik ekran görüntüsü."""
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)

    yol = _ekran_yolu(d["dosya"], d.get("_dil", "tr"))
    _yazi(slayt, 0.62, 0.98, 8.6, 0.5, [(d["baslik"], 26, Z["yazi_ac"], True, 0)])
    if d.get("alt"):
        _yazi(slayt, 0.62, 1.50, 11.0, 0.35, [(d["alt"], 12, Z["sonuk"], False, 0)])

    # Genislik 16:9 yuksekligi belirler; slayt 7.5 inc oldugu icin gorsel
    # 1.95'ten baslayinca en fazla ~5.3 inc yuksek olabilir -> 9.4 inc genislik.
    # Ilk surumde 10.7 verilmisti ve goruntu slaytin altindan tasiyordu.
    g = 9.4
    x = (13.333 - g) / 2
    if yol:
        _gorsel(slayt, yol, x, 1.95, g)
    else:
        _kutu(slayt, x, 1.95, g, g * 9.0 / 16.0, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])


def _s_galeri(prs, slayt, d, surum, marka, sayfa, toplam):
    """3x2 küçük görsel ızgarası — geri kalan ekranları tek slaytta gösterir."""
    _zemin(slayt)
    _ustbilgi(slayt, marka, surum, sayfa, toplam)
    _slayt_basligi(slayt, d["baslik"], d.get("alt"))

    # Olculer 2 satir + 2 etiket satirinin 7.5 inclik slayta SIGMASI icin
    # geriye dogru hesaplandi. Ilk surumde genislik 3.81 incti ve alt satirin
    # etiketleri slaytin disinda kaliyordu.
    ogeler = d["ogeler"]
    sutun = 3
    w, h = 3.52, 3.52 * 9.0 / 16.0        # 1.98 inc yukseklik
    bosluk_x, bosluk_y = 0.34, 0.45
    ust = 2.30
    sol = (13.333 - (sutun * w + (sutun - 1) * bosluk_x)) / 2
    for i, (dosya, etiket) in enumerate(ogeler[:6]):
        cx = sol + (i % sutun) * (w + bosluk_x)
        cy = ust + (i // sutun) * (h + bosluk_y + 0.28)
        yol = _ekran_yolu(dosya, d.get("_dil", "tr"))
        if yol:
            _gorsel(slayt, yol, cx, cy, w)
        else:
            _kutu(slayt, cx, cy, w, h, dolgu=Z["kart"], cizgi=Z["kart_cizgi"])
        _yazi(slayt, cx, cy + h + 0.08, w, 0.28,
              [(etiket, 10, Z["sonuk"], False, 0)], hiza="c")


CIZERLER = {
    "kapak": _s_kapak, "bolum": _s_bolum, "vurgu": _s_vurgu,
    "kartlar": _s_kartlar, "sayilar": _s_sayilar, "zincir": _s_zincir,
    "iki_sutun": _s_iki_sutun, "kapanis": _s_kapanis,
    "ekran": _s_ekran, "ekran_tam": _s_ekran_tam, "galeri": _s_galeri,
}


# ══════════════════════════════════════════════════════════════════════════
# PPTX ÜRETİMİ
# ══════════════════════════════════════════════════════════════════════════
def pptx_uret(dil, surum):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)      # 16:9
    prs.slide_height = Inches(7.5)

    tanimlar = SLAYTLAR[dil]
    marka = METIN[dil]["marka"]
    toplam = len(tanimlar)
    bos = prs.slide_layouts[6]            # tamamen boş düzen

    for i, d in enumerate(tanimlar, 1):
        slayt = prs.slides.add_slide(bos)
        # Ekran slaytları hangi dilin görüntüsünü alacağını bilmeli.
        d = dict(d, _dil=dil)
        CIZERLER[d["tip"]](prs, slayt, d, surum, marka, i, toplam)

    yol = _yollar(dil)[0]
    os.makedirs(os.path.dirname(yol), exist_ok=True)
    prs.save(yol)
    return yol, toplam


# ══════════════════════════════════════════════════════════════════════════
# BASKI SÜRÜMÜ
# ══════════════════════════════════════════════════════════════════════════
def _sekil_baskiya(sekil):
    """Bir şeklin dolgu, çizgi ve yazı renklerini baskı paletine çevirir."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if sekil.shape_type == MSO_SHAPE_TYPE.GROUP:
        for alt in sekil.shapes:
            _sekil_baskiya(alt)
        return

    try:
        f = sekil.fill
        if f.type is not None and f.type == 1:      # solid
            yeni = BASKI_ZEMIN.get(str(f.fore_color.rgb).upper())
            if yeni:
                f.fore_color.rgb = _renk(yeni)
    except Exception:
        pass

    try:
        ln = sekil.line
        if ln.fill.type == 1:
            yeni = BASKI_ZEMIN.get(str(ln.color.rgb).upper())
            if yeni:
                ln.color.rgb = _renk(yeni)
    except Exception:
        pass

    if sekil.has_text_frame:
        for p in sekil.text_frame.paragraphs:
            for r in p.runs:
                try:
                    renk = r.font.color
                    if renk and renk.type is not None:
                        yeni = BASKI_YAZI.get(str(renk.rgb).upper())
                        if yeni:
                            renk.rgb = _renk(yeni)
                except Exception:
                    pass


def baskiya_cevir(prs):
    """Sunumun tamamını baskı paletine çevirir → değiştirilen şekil sayısı."""
    sayac = 0
    for slayt in prs.slides:
        try:
            arka = slayt.background.fill
            if arka.type == 1:
                yeni = BASKI_ZEMIN.get(str(arka.fore_color.rgb).upper())
                if yeni:
                    arka.fore_color.rgb = _renk(yeni)
        except Exception:
            pass
        for sekil in slayt.shapes:
            _sekil_baskiya(sekil)
            sayac += 1
    return sayac


# ══════════════════════════════════════════════════════════════════════════
# POWERPOINT COM — PDF ve PNG
# ══════════════════════════════════════════════════════════════════════════
def powerpoint_disa_aktar(pptx, pdf, png_dizin=None):
    """PowerPoint ile PDF (ve istenirse slayt PNG'leri) üretir.

    PowerPoint yoksa (False, sebep) döner; çağıran karar verir. PNG üretimi
    yalnızca HTML'i olan sürümler için istenir — 22 slaytta saniyeler sürüyor.
    """
    try:
        import win32com.client
    except ImportError:
        return False, "pywin32 kurulu değil"

    uygulama = None
    sunum = None
    try:
        uygulama = win32com.client.Dispatch("PowerPoint.Application")
    except Exception as exc:
        return False, "PowerPoint açılamadı (%s)" % exc

    try:
        sunum = uygulama.Presentations.Open(
            os.path.abspath(pptx), WithWindow=False, ReadOnly=False)
        sunum.SaveAs(os.path.abspath(pdf), 32)          # 32 = ppSaveAsPDF
        if png_dizin:
            os.makedirs(png_dizin, exist_ok=True)
            sunum.Export(os.path.abspath(png_dizin), "PNG", 1920, 1080)
        return True, None
    except Exception as exc:
        return False, str(exc)
    finally:
        try:
            if sunum is not None:
                sunum.Close()
        except Exception:
            pass
        try:
            if uygulama is not None:
                uygulama.Quit()
        except Exception:
            pass


# ══════════════════════════════════════════════════════════════════════════
# HTML — slayt PNG'leri gömülü, tek dosya
# ══════════════════════════════════════════════════════════════════════════
_HTML_KALIP = """<!doctype html>
<html lang="__DIL__">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1, viewport-fit=cover">
<meta name="theme-color" content="#0F172A">
<title>__BASLIK__</title>
<style>
  *{box-sizing:border-box;margin:0;padding:0}
  html,body{height:100%;background:#0F172A;color:#F1F5F9;
    font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,Calibri,sans-serif;
    overscroll-behavior:none;-webkit-text-size-adjust:100%}
  body{display:flex;flex-direction:column;height:100dvh}
  header{flex:0 0 auto;display:flex;align-items:center;gap:.6rem;
    padding:calc(.5rem + env(safe-area-inset-top)) .9rem .5rem;
    border-bottom:1px solid #1E293B;background:#0F172A}
  .brand{font-size:.86rem;font-weight:600;letter-spacing:.02em;white-space:nowrap;
    overflow:hidden;text-overflow:ellipsis}
  .ver{flex:0 0 auto;font-size:.7rem;font-weight:700;color:#0F172A;background:#2DD4BF;
    border-radius:99px;padding:.15rem .5rem}
  .count{margin-left:auto;flex:0 0 auto;font-size:.8rem;color:#94A3B8;
    font-variant-numeric:tabular-nums}
  .bar{flex:0 0 auto;height:2px;background:#1E293B}
  .bar > i{display:block;height:100%;width:0;background:#4F46E5;transition:width .18s ease}
  main{flex:1 1 auto;display:flex;overflow-x:auto;overflow-y:hidden;
    scroll-snap-type:x mandatory;-webkit-overflow-scrolling:touch;
    scrollbar-width:none}
  main::-webkit-scrollbar{display:none}
  .slide{flex:0 0 100%;scroll-snap-align:center;display:flex;align-items:center;
    justify-content:center;padding:.6rem}
  .slide img{max-width:100%;max-height:100%;object-fit:contain;border-radius:10px;
    box-shadow:0 10px 34px -12px rgba(0,0,0,.7)}
  footer{flex:0 0 auto;display:flex;align-items:center;gap:.5rem;
    padding:.5rem .9rem calc(.5rem + env(safe-area-inset-bottom));
    border-top:1px solid #1E293B;background:#0F172A}
  button{appearance:none;border:1px solid #334155;background:#1E293B;color:#F1F5F9;
    border-radius:8px;padding:.45rem .8rem;font-size:.82rem;font-weight:500;
    cursor:pointer;font-family:inherit}
  button:active{background:#334155}
  button:disabled{opacity:.4;cursor:default}
  .hint{margin:0 auto;font-size:.72rem;color:#64748B;text-align:center}
  @media (min-width:900px){.hint{font-size:.78rem}}
  noscript{display:block;padding:2rem;text-align:center;color:#94A3B8}
</style>
</head>
<body>
<header>
  <span class="brand">__MARKA__</span>
  <span class="ver">__SURUM__</span>
  <span class="count"><b id="cur">1</b> / __ADET__</span>
</header>
<div class="bar"><i id="prog"></i></div>
<main id="deck">
__SLAYTLAR__
</main>
<footer>
  <button id="prev" aria-label="__ONCEKI__">‹ __ONCEKI__</button>
  <span class="hint">__IPUCU__</span>
  <button id="next" aria-label="__SONRAKI__">__SONRAKI__ ›</button>
</footer>
<noscript>__UYARI__</noscript>
<script>
(function(){
  var deck=document.getElementById('deck'),cur=document.getElementById('cur'),
      prog=document.getElementById('prog'),prev=document.getElementById('prev'),
      next=document.getElementById('next'),n=__ADET__;
  function idx(){return Math.round(deck.scrollLeft/deck.clientWidth);}
  function paint(){var i=idx();cur.textContent=i+1;
    prog.style.width=((i+1)/n*100)+'%';
    prev.disabled=i<=0;next.disabled=i>=n-1;}
  function go(i){i=Math.max(0,Math.min(n-1,i));
    deck.scrollTo({left:i*deck.clientWidth,behavior:'smooth'});}
  deck.addEventListener('scroll',function(){
    clearTimeout(deck._t);deck._t=setTimeout(paint,60);},{passive:true});
  prev.addEventListener('click',function(){go(idx()-1);});
  next.addEventListener('click',function(){go(idx()+1);});
  document.addEventListener('keydown',function(e){
    if(e.key==='ArrowRight'||e.key===' '||e.key==='PageDown'){e.preventDefault();go(idx()+1);}
    else if(e.key==='ArrowLeft'||e.key==='PageUp'){e.preventDefault();go(idx()-1);}
    else if(e.key==='Home'){e.preventDefault();go(0);}
    else if(e.key==='End'){e.preventDefault();go(n-1);}});
  window.addEventListener('resize',function(){go(idx());},{passive:true});
  paint();
})();
</script>
</body>
</html>
"""


def html_uret(pngler, dil, surum):
    """Slayt PNG'lerini tek dosyalık mobil sunuma gömer → yazılan bayt sayısı."""
    yol = _yollar(dil)[2]
    M = METIN[dil]
    parcalar = []
    for i, png in enumerate(pngler, 1):
        with open(png, "rb") as f:
            veri = base64.b64encode(f.read()).decode("ascii")
        # İlk iki slayt hemen, gerisi tembel yüklenir: 22 slaytlık base64 yükünü
        # tek seferde çözmek mobilde ilk açılışı saniyelerce geciktiriyor.
        yukleme = "eager" if i <= 2 else "lazy"
        parcalar.append(
            '<figure class="slide" id="s%d"><img src="data:image/png;base64,%s" '
            'alt="%s %d" loading="%s" decoding="async"></figure>'
            % (i, veri, M["slayt"], i, yukleme))

    html = _HTML_KALIP
    for anahtar, deger in (
            ("__DIL__", M["kod"]), ("__BASLIK__", M["baslik"]),
            ("__MARKA__", M["marka"]), ("__ONCEKI__", M["onceki"]),
            ("__SONRAKI__", M["sonraki"]), ("__IPUCU__", M["ipucu"]),
            ("__UYARI__", M["uyari"]), ("__SURUM__", surum),
            ("__ADET__", str(len(pngler))), ("__SLAYTLAR__", "\n".join(parcalar))):
        html = html.replace(anahtar, deger)

    with open(yol, "w", encoding="utf-8") as f:
        f.write(html)
    return yol, len(html.encode("utf-8"))


# ══════════════════════════════════════════════════════════════════════════
# AKIŞ
# ══════════════════════════════════════════════════════════════════════════
def uret(dil, sadece_pptx=False):
    surum = _surum()[0]
    etiket = "TR" if dil == "tr" else "EN"

    pptx, pdf, html = _yollar(dil)
    b_pptx, b_pdf = _baski_yollari(dil)

    yol, adet = pptx_uret(dil, surum)
    print("   %s [%s] PPTX  : %d slayt  (%.0f KB)"
          % (_iz("✓"), etiket, adet, os.path.getsize(yol) / 1024))

    if sadece_pptx:
        return True

    # ── Baskı sürümü: aynı slaytlar, beyaz zemin ─────────────────────────
    from pptx import Presentation
    sunum = Presentation(pptx)
    degisen = baskiya_cevir(sunum)
    sunum.save(b_pptx)
    print("   %s [%s] BASKI : %d şekil beyaz zemine çevrildi  (%.0f KB)"
          % (_iz("✓"), etiket, degisen, os.path.getsize(b_pptx) / 1024))

    # ── PDF + PNG ────────────────────────────────────────────────────────
    png_dizin = tempfile.mkdtemp(prefix="vs_slayt_")
    try:
        ok, hata = powerpoint_disa_aktar(pptx, pdf, png_dizin)
        if not ok:
            print("   %s [%s] PDF/PNG atlandı: %s" % (_iz("✗"), etiket, hata))
            print("        (PPTX üretildi; PDF ve HTML için PowerPoint gerekir)")
            return True
        print("   %s [%s] PDF   : %.1f MB"
              % (_iz("✓"), etiket, os.path.getsize(pdf) / 1048576))

        ok2, hata2 = powerpoint_disa_aktar(b_pptx, b_pdf)     # baskıda PNG yok
        if ok2:
            print("   %s [%s] BASKI PDF : %.1f MB"
                  % (_iz("✓"), etiket, os.path.getsize(b_pdf) / 1048576))
        else:
            print("   %s [%s] Baskı PDF atlandı: %s" % (_iz("✗"), etiket, hata2))

        # Windows dosya sistemi harf duyarsız olduğu icin "*.PNG" ve "*.png"
        # AYNI dosyalari dondurur; ikisini toplamak her slaydi iki kez gomer.
        # Tek desenle tarayip uzantiyi kendimiz suzuyoruz.
        pngler = sorted(
            (y for y in glob.glob(os.path.join(png_dizin, "*"))
             if y.lower().endswith(".png")),
            key=lambda y: int(re.findall(r"(\d+)", os.path.basename(y))[-1]))
        if not pngler:
            print("   %s [%s] PNG bulunamadı, HTML atlandı" % (_iz("✗"), etiket))
            return True

        hyol, bayt = html_uret(pngler, dil, surum)
        print("   %s [%s] HTML  : %d slayt gömülü  (%.1f MB)"
              % (_iz("✓"), etiket, len(pngler), bayt / 1048576))
    finally:
        shutil.rmtree(png_dizin, ignore_errors=True)

    return True


def durum():
    surum = _surum()
    print("   Sürüm: %s  (kaynak: %s)" % surum)
    for dil in DILLER:
        etiket = "TR" if dil == "tr" else "EN"
        for ad, yol in zip(("PPTX", "PDF ", "HTML"), _yollar(dil)):
            if os.path.exists(yol):
                print("   [%s] %s : %.2f MB   %s"
                      % (etiket, ad, os.path.getsize(yol) / 1048576,
                         os.path.basename(yol)))
            else:
                print("   [%s] %s : yok" % (etiket, ad))
        for ad, yol in zip(("BASKI PPTX", "BASKI PDF "), _baski_yollari(dil)):
            if os.path.exists(yol):
                print("   [%s] %s : %.2f MB   %s"
                      % (etiket, ad, os.path.getsize(yol) / 1048576,
                         os.path.basename(yol)))


def main():
    ap = argparse.ArgumentParser(description="Van Sales tanıtım sunumu üretici")
    ap.add_argument("--dil", choices=["tr", "en", "hepsi"], default="hepsi")
    ap.add_argument("--sadece-pptx", action="store_true",
                    help="PowerPoint'e dokunma, yalnız PPTX üret")
    ap.add_argument("--kontrol", action="store_true",
                    help="üretmeden mevcut çıktıları raporla")
    a = ap.parse_args()

    if a.kontrol:
        durum()
        return 0

    os.makedirs(CIKTI, exist_ok=True)
    surum, kaynak = _surum()
    print("Van Sales tanıtım sunumu üretiliyor  %s  (sürüm kaynağı: %s)"
          % (surum, kaynak))
    print("Çıktı dizini: %s\n" % CIKTI)

    diller = DILLER if a.dil == "hepsi" else (a.dil,)
    for dil in diller:
        uret(dil, sadece_pptx=a.sadece_pptx)
        print("")

    durum()
    return 0


if __name__ == "__main__":
    sys.exit(main())
